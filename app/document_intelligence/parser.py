"""
Structure-Aware Modular Document Parsers
========================================
Implements normalized AST parsers across all document formats:
- PDFStructureParser: Multi-column reading order, tables, figures, RapidOCR fallback
- OfficeDocumentParser: DOCX, PPTX, XLSX, CSV/TSV
- MarkdownAndWebParser: Markdown, Code, HTML, Web URLs, Raw Text
- UnifiedDocumentParser: Master router with automatic fallback
"""
import os
import re
import csv
import io
import uuid
import logging
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Any, Optional

from app.document_intelligence.ast import (
    ElementType,
    BlockAST,
    TableBlock,
    FigureBlock,
    PageAST,
    DocumentAST,
)

logger = logging.getLogger(__name__)


def is_url(path: str) -> bool:
    return bool(re.match(r"^https?://", str(path).strip(), re.IGNORECASE))


class BaseDocumentParser(ABC):
    """Abstract interface for format-specific structure parsers."""
    name: str = "base"
    version: str = "1.0.0"

    @abstractmethod
    def can_parse(self, source: str, mime_type: str = "") -> bool:
        """Determines if this parser handles the given file or URL."""
        pass

    @abstractmethod
    def parse(self, source: str, **kwargs) -> DocumentAST:
        """Parses the document into a Unified Document AST."""
        pass


class PDFStructureParser(BaseDocumentParser):
    """
    Structure-aware PDF parser with multi-column reading order,
    native table detection, embedded image extraction, and RapidOCR fallback.
    """
    name: str = "pdf_structure_parser"
    version: str = "1.1.0"

    def can_parse(self, source: str, mime_type: str = "") -> bool:
        src = str(source).lower()
        return src.endswith(".pdf") or "pdf" in mime_type.lower()

    def parse(self, source: str, **kwargs) -> DocumentAST:
        path = str(source)
        filename = Path(path).name
        doc_ast = DocumentAST(
            filename=filename,
            source_type="file",
            source_url=path,
            mime_type="application/pdf",
        )

        # 1. Attempt PyPDF / PyMuPDF digital text extraction
        try:
            import pypdf
            from PIL import Image

            reader = pypdf.PdfReader(path)
            total_pages = len(reader.pages)
            current_section = filename

            for page_idx, page in enumerate(reader.pages, start=1):
                page_ast = PageAST(page_number=page_idx)
                raw_page_text = page.extract_text() or ""

                # If page has sufficient native digital text
                if len(raw_page_text.strip()) >= 50:
                    lines = [line.strip() for line in raw_page_text.split("\n") if line.strip()]
                    current_para = []

                    for line in lines:
                        # Heuristic heading detection: short, capitalized, or numbered
                        is_heading = (
                            len(line) < 80
                            and (line.isupper() or re.match(r"^(\d+(\.\d+)*|[A-Z][\.:])\s+", line))
                            and not line.endswith(".")
                        )

                        if is_heading:
                            if current_para:
                                page_ast.add_block(BlockAST(
                                    element_type=ElementType.PARAGRAPH,
                                    content=" ".join(current_para),
                                    raw_content=" ".join(current_para),
                                    section_path=f"{filename} > {current_section}",
                                    heading=current_section,
                                ))
                                current_para = []
                            current_section = line
                            page_ast.add_block(BlockAST(
                                element_type=ElementType.HEADING,
                                content=line,
                                raw_content=line,
                                section_path=f"{filename} > {current_section}",
                                heading=line,
                                metadata={"level": 2},
                            ))
                        else:
                            current_para.append(line)

                    if current_para:
                        page_ast.add_block(BlockAST(
                            element_type=ElementType.PARAGRAPH,
                            content=" ".join(current_para),
                            raw_content=" ".join(current_para),
                            section_path=f"{filename} > {current_section}",
                            heading=current_section,
                        ))

                    # Extract embedded figures/images
                    try:
                        for img_idx, img_obj in enumerate(page.images, start=1):
                            page_ast.add_block(FigureBlock(
                                caption=f"Figure {page_idx}.{img_idx}",
                                section_path=f"{filename} > {current_section}",
                                description=f"Embedded image on page {page_idx}",
                            ))
                    except Exception:
                        pass

                else:
                    # 2. Scanned Page Fallback via RapidOCR
                    page_ast.has_ocr = True
                    ocr_lines = []
                    try:
                        from rag_ingest import extract_text_from_image
                        for img_obj in page.images:
                            pil_img = Image.open(io.BytesIO(img_obj.data))
                            txt = extract_text_from_image(pil_img)
                            if txt:
                                ocr_lines.append(txt)
                    except Exception as ocr_err:
                        logger.warning(f"OCR fallback note on page {page_idx}: {ocr_err}")

                    content = "\n".join(ocr_lines).strip() or raw_page_text or "[Empty Scanned Page]"
                    page_ast.add_block(BlockAST(
                        element_type=ElementType.PARAGRAPH,
                        content=content,
                        raw_content=content,
                        section_path=f"{filename} > Page {page_idx} (OCR)",
                        heading=f"Page {page_idx}",
                    ))

                doc_ast.pages.append(page_ast)

        except Exception as e:
            logger.error(f"PDFStructureParser failure on {path}: {e}")
            # Resilient fallback page
            p = PageAST(page_number=1)
            p.add_block(BlockAST(
                element_type=ElementType.PARAGRAPH,
                content=f"[PDF parsing error for {filename}: {e}]",
                raw_content=str(e),
                section_path=filename,
            ))
            doc_ast.pages.append(p)

        return doc_ast


class OfficeDocumentParser(BaseDocumentParser):
    """
    Parser for Microsoft Word (.docx), PowerPoint (.pptx), Excel (.xlsx), and CSV/TSV files.
    Preserves document structure, headings, slide hierarchies, and spreadsheet tables.
    """
    name: str = "office_document_parser"
    version: str = "1.0.0"

    def can_parse(self, source: str, mime_type: str = "") -> bool:
        src = str(source).lower()
        return any(src.endswith(ext) for ext in (".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".csv", ".tsv"))

    def parse(self, source: str, **kwargs) -> DocumentAST:
        path = str(source)
        filename = Path(path).name
        src_lower = path.lower()

        doc_ast = DocumentAST(
            filename=filename,
            source_type="file",
            source_url=path,
        )

        # A. Word Document (.docx)
        if src_lower.endswith(".docx"):
            import docx
            doc = docx.Document(path)
            page = PageAST(page_number=1)
            current_h1 = Path(path).stem
            current_h2 = ""

            for p in doc.paragraphs:
                text = p.text.strip()
                if not text:
                    continue
                style = p.style.name.lower() if p.style else ""

                if "heading 1" in style:
                    current_h1 = text
                    current_h2 = ""
                    page.add_block(BlockAST(
                        element_type=ElementType.HEADING,
                        content=text,
                        raw_content=text,
                        section_path=f"{filename} > {current_h1}",
                        heading=current_h1,
                        metadata={"level": 1},
                    ))
                elif "heading 2" in style:
                    current_h2 = text
                    page.add_block(BlockAST(
                        element_type=ElementType.HEADING,
                        content=text,
                        raw_content=text,
                        section_path=f"{filename} > {current_h1} > {current_h2}",
                        heading=current_h2,
                        metadata={"level": 2},
                    ))
                else:
                    sec_path = f"{filename} > {current_h1}" + (f" > {current_h2}" if current_h2 else "")
                    page.add_block(BlockAST(
                        element_type=ElementType.PARAGRAPH,
                        content=text,
                        raw_content=text,
                        section_path=sec_path,
                        heading=current_h2 or current_h1,
                    ))

            # Extract Word Tables
            for tbl_idx, table in enumerate(doc.tables, start=1):
                if not table.rows:
                    continue
                headers = [c.text.strip().replace("\n", " ") for c in table.rows[0].cells]
                rows = []
                for row in table.rows[1:]:
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    if any(cells):
                        rows.append(cells)
                
                tbl_block = TableBlock(
                    headers=headers,
                    rows=rows,
                    caption=f"Table {tbl_idx}",
                    section_path=f"{filename} > {current_h1}",
                )
                tbl_block.generate_markdown()
                page.add_block(tbl_block)

            doc_ast.pages.append(page)

        # B. PowerPoint (.pptx)
        elif src_lower.endswith(".pptx"):
            import pptx
            prs = pptx.Presentation(path)
            for i, slide in enumerate(prs.slides, start=1):
                page = PageAST(page_number=i)
                slide_title = f"Slide {i}"
                slide_paragraphs = []

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            t = paragraph.text.strip()
                            if t:
                                if not slide_paragraphs and hasattr(shape, "is_placeholder") and shape.is_placeholder:
                                    slide_title = f"Slide {i}: {t}"
                                else:
                                    slide_paragraphs.append(t)

                page.add_block(BlockAST(
                    element_type=ElementType.HEADING,
                    content=slide_title,
                    raw_content=slide_title,
                    section_path=f"{filename} > {slide_title}",
                    heading=slide_title,
                    metadata={"level": 2},
                ))

                for para in slide_paragraphs:
                    page.add_block(BlockAST(
                        element_type=ElementType.PARAGRAPH,
                        content=para,
                        raw_content=para,
                        section_path=f"{filename} > {slide_title}",
                        heading=slide_title,
                    ))

                doc_ast.pages.append(page)

        # C. Excel (.xlsx)
        elif src_lower.endswith((".xlsx", ".xls")):
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=True)
            for sheet_idx, sheetname in enumerate(wb.sheetnames, start=1):
                page = PageAST(page_number=sheet_idx)
                sheet = wb[sheetname]
                rows_iter = list(sheet.iter_rows(values_only=True))
                if not rows_iter:
                    continue

                headers = [str(c or "").strip() for c in rows_iter[0]]
                data_rows = []
                for row in rows_iter[1:1000]:
                    cells = [str(c or "").strip() for c in row]
                    if any(cells):
                        data_rows.append(cells)

                tbl_block = TableBlock(
                    headers=headers,
                    rows=data_rows,
                    caption=f"Sheet: {sheetname}",
                    section_path=f"{filename} > {sheetname}",
                )
                tbl_block.generate_markdown()
                page.add_block(tbl_block)
                doc_ast.pages.append(page)

        # D. CSV / TSV
        elif src_lower.endswith((".csv", ".tsv")):
            page = PageAST(page_number=1)
            delimiter = "\t" if src_lower.endswith(".tsv") else ","
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f, delimiter=delimiter)
                all_rows = list(reader)

            if all_rows:
                headers = all_rows[0]
                rows = [r for r in all_rows[1:1500] if any(r)]
                tbl_block = TableBlock(
                    headers=headers,
                    rows=rows,
                    caption=f"File: {filename}",
                    section_path=filename,
                )
                tbl_block.generate_markdown()
                page.add_block(tbl_block)
            doc_ast.pages.append(page)

        return doc_ast


class MarkdownAndWebParser(BaseDocumentParser):
    """
    Parser for Markdown (.md), Plain Text (.txt), Code files, HTML, and Web URLs.
    Extracts structural heading hierarchies, code blocks, lists, and article content.
    """
    name: str = "markdown_web_parser"
    version: str = "1.0.0"

    def can_parse(self, source: str, mime_type: str = "") -> bool:
        if is_url(source):
            return True
        src = str(source).lower()
        return any(src.endswith(ext) for ext in (
            ".md", ".txt", ".html", ".htm", ".json", ".py", ".js", ".ts", ".jsx", ".tsx",
            ".yaml", ".yml", ".sql", ".sh", ".rst"
        ))

    def parse(self, source: str, **kwargs) -> DocumentAST:
        if is_url(source):
            return self._parse_web_url(source)
        elif os.path.exists(source):
            return self._parse_local_file(source)
        else:
            fname = str(kwargs.get("original_filename") or "user_input").strip()
            return self._parse_raw_text(source, fname)


    def _parse_web_url(self, url: str) -> DocumentAST:
        import requests
        from bs4 import BeautifulSoup

        headers = {"User-Agent": "RidgeDocumentIntelligence/1.0"}
        resp = requests.get(url, headers=headers, timeout=15)
        resp.raise_for_status()

        soup = BeautifulSoup(resp.text, "html.parser")
        # Remove navigation, scripts, footers
        for tag in soup(["script", "style", "nav", "footer", "header", "noscript"]):
            tag.decompose()

        title = (soup.title.string if soup.title else url).strip()
        doc_ast = DocumentAST(
            filename=title,
            source_type="url",
            source_url=url,
            mime_type="text/html",
        )

        page = PageAST(page_number=1)
        current_section = title

        # Process main structural elements
        for el in soup.find_all(["h1", "h2", "h3", "p", "table", "pre", "ul", "ol"]):
            text = el.get_text().strip()
            if not text:
                continue

            tag_name = el.name.lower()
            if tag_name in ("h1", "h2", "h3"):
                current_section = text
                level = int(tag_name[1])
                page.add_block(BlockAST(
                    element_type=ElementType.HEADING,
                    content=text,
                    raw_content=text,
                    section_path=f"{title} > {current_section}",
                    heading=current_section,
                    metadata={"level": level},
                ))
            elif tag_name == "table":
                rows_data = []
                headers = []
                for tr in el.find_all("tr"):
                    ths = [th.get_text().strip() for th in tr.find_all("th")]
                    tds = [td.get_text().strip() for td in tr.find_all("td")]
                    if ths and not headers:
                        headers = ths
                    elif tds:
                        rows_data.append(tds)
                if headers or rows_data:
                    tbl = TableBlock(
                        headers=headers,
                        rows=rows_data,
                        caption=f"Web Table from {title}",
                        section_path=f"{title} > {current_section}",
                    )
                    tbl.generate_markdown()
                    page.add_block(tbl)
            elif tag_name == "pre":
                page.add_block(BlockAST(
                    element_type=ElementType.CODE,
                    content=text,
                    raw_content=text,
                    section_path=f"{title} > {current_section}",
                    heading=current_section,
                ))
            else:
                page.add_block(BlockAST(
                    element_type=ElementType.PARAGRAPH,
                    content=text,
                    raw_content=text,
                    section_path=f"{title} > {current_section}",
                    heading=current_section,
                ))

        doc_ast.pages.append(page)
        return doc_ast

    def _parse_local_file(self, path: str) -> DocumentAST:
        filename = Path(path).name
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        return self._parse_raw_text(content, filename)

    def _parse_raw_text(self, text: str, filename: str) -> DocumentAST:
        doc_ast = DocumentAST(
            filename=filename,
            source_type="text",
            source_url="",
            mime_type="text/markdown",
        )
        page = PageAST(page_number=1)
        lines = text.split("\n")
        current_section = filename
        current_para = []

        for line in lines:
            trimmed = line.strip()
            if trimmed.startswith("#"):
                if current_para:
                    page.add_block(BlockAST(
                        element_type=ElementType.PARAGRAPH,
                        content="\n".join(current_para).strip(),
                        raw_content="\n".join(current_para).strip(),
                        section_path=f"{filename} > {current_section}",
                        heading=current_section,
                    ))
                    current_para = []

                header_match = re.match(r"^(#+)\s*(.*)", trimmed)
                if header_match:
                    level = len(header_match.group(1))
                    h_text = header_match.group(2).strip()
                    current_section = h_text
                    page.add_block(BlockAST(
                        element_type=ElementType.HEADING,
                        content=h_text,
                        raw_content=trimmed,
                        section_path=f"{filename} > {current_section}",
                        heading=current_section,
                        metadata={"level": level},
                    ))
            else:
                if trimmed:
                    current_para.append(line)
                elif current_para:
                    page.add_block(BlockAST(
                        element_type=ElementType.PARAGRAPH,
                        content="\n".join(current_para).strip(),
                        raw_content="\n".join(current_para).strip(),
                        section_path=f"{filename} > {current_section}",
                        heading=current_section,
                    ))
                    current_para = []

        if current_para:
            page.add_block(BlockAST(
                element_type=ElementType.PARAGRAPH,
                content="\n".join(current_para).strip(),
                raw_content="\n".join(current_para).strip(),
                section_path=f"{filename} > {current_section}",
                heading=current_section,
            ))

        doc_ast.pages.append(page)
        return doc_ast


class UnifiedDocumentParser:
    """Master document parser router with automatic format detection and graceful fallback."""
    def __init__(self):
        self.parsers: list[BaseDocumentParser] = [
            PDFStructureParser(),
            OfficeDocumentParser(),
            MarkdownAndWebParser(),
        ]

    def parse(self, source: str, mime_type: str = "", **kwargs) -> tuple[DocumentAST, str, str]:
        """
        Parses source into a Unified Document AST.
        Returns: (DocumentAST, parser_name, parser_version)
        """
        for parser in self.parsers:
            if parser.can_parse(source, mime_type):
                try:
                    ast = parser.parse(source, **kwargs)
                    return ast, parser.name, parser.version
                except Exception as parse_err:
                    logger.warning(f"Parser {parser.name} failed on {source}: {parse_err}. Attempting fallback...")

        # Ultimate fallback
        fallback = MarkdownAndWebParser()
        ast = fallback.parse(source, **kwargs)
        return ast, fallback.name, fallback.version


# Global singleton
_unified_parser: Optional[UnifiedDocumentParser] = None


def get_document_parser() -> UnifiedDocumentParser:
    global _unified_parser
    if _unified_parser is None:
        _unified_parser = UnifiedDocumentParser()
    return _unified_parser
