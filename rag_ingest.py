"""
Ridge: Universal Ingestion & Knowledge Crag Engine
===================================================
Supports multi-format parsing, hierarchical chunking, and persistent embedding into ChromaDB.

Supported Sources:
- Documents: PDF (.pdf), Microsoft Word (.docx), PowerPoint (.pptx), Markdown (.md), Plain Text (.txt)
- Tabular Data: Excel (.xlsx), CSV (.csv), TSV (.tsv)
- Code & Config: Python (.py), JavaScript/TypeScript (.js, .jsx, .ts, .tsx), Web (.html, .css),
                 Config (.json, .yaml, .yml, .toml, .ini), Database (.sql), Systems (.c, .cpp, .java, .go, .rs, .sh)
- Media & Subtitles: SubRip (.srt), WebVTT (.vtt)
- Web & Media URLs: YouTube Transcripts, GitHub Repos/Files, ArXiv Papers, and standard Web URLs.
"""

import os
import re
import csv
from pathlib import Path
from urllib.parse import urlparse, parse_qs
from typing import Optional

import requests
from bs4 import BeautifulSoup
from langchain_core.documents import Document
from langchain_text_splitters import MarkdownHeaderTextSplitter, RecursiveCharacterTextSplitter

DEFAULT_ARTICLE_URL = "https://lilianweng.github.io/posts/2023-06-23-agent/"
DEFAULT_CHUNK_SIZE = 1500
DEFAULT_CHUNK_OVERLAP = 200
DEFAULT_USER_AGENT = "RidgeCRAG/1.0"


# ---------------------------------------------------------------------------
# Shared Chunking and Heading Preservation Helpers
# ---------------------------------------------------------------------------

def clean_heading(text: str) -> str:
    return text.strip().removesuffix("#").strip()


def heading_lines(metadata: dict) -> list[str]:
    return [metadata[key] for key in ("h1", "h2", "h3") if metadata.get(key)]


def breadcrumb_string(metadata: dict) -> str:
    parts = []
    source = metadata.get("source")
    if source:
        src_name = source.split("/")[-1] if not is_url(source) else source
        parts.append(src_name)
    for key in ("h1", "h2", "h3"):
        val = metadata.get(key)
        if val:
            clean = clean_heading(str(val))
            if clean and clean not in parts:
                parts.append(clean)
    return " > ".join(parts) if parts else ""


def format_section(section: Document, source: str) -> Document:
    metadata = {"source": source, **section.metadata}
    breadcrumb = breadcrumb_string(metadata)
    if breadcrumb:
        metadata["breadcrumb"] = breadcrumb
    page_content = "\n".join(heading_lines(section.metadata) + [section.page_content]).strip()
    return Document(page_content=page_content, metadata=metadata)


def ensure_chunk_has_headings(chunk: Document) -> Document:
    meta = dict(chunk.metadata)
    breadcrumb = breadcrumb_string(meta)
    if breadcrumb:
        meta["breadcrumb"] = breadcrumb
    
    headings = heading_lines(meta)
    prefix_header = f"[Context: {breadcrumb}]" if breadcrumb else ""
    
    content = chunk.page_content
    if prefix_header and not content.startswith("[Context:"):
        content = f"{prefix_header}\n{content}"
    elif headings and not any(content.startswith(h) for h in headings):
        content = "\n".join(headings + [content]).strip()

    return Document(page_content=content, metadata=meta)


def _sub_chunk(sections: list[Document], size: int, overlap: int) -> list[Document]:
    """Runs long sections through a character-based splitter, preserving hierarchy."""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=size,
        chunk_overlap=overlap,
        separators=["\n\n", "\n", "```", ".", " ", ""],
    )
    return [ensure_chunk_has_headings(chunk) for chunk in text_splitter.split_documents(sections)]


# ---------------------------------------------------------------------------
# Semantic Chunking by Embedding Gradient
# ---------------------------------------------------------------------------

def _cosine_sim(a: list[float], b: list[float]) -> float:
    """Fast cosine similarity between two embedding vectors."""
    dot = sum(x * y for x, y in zip(a, b))
    na = sum(x * x for x in a) ** 0.5
    nb = sum(x * x for x in b) ** 0.5
    return dot / (na * nb + 1e-9)


def semantic_chunk(
    text: str,
    embedder,
    window: int = 3,
    percentile: float = 25.0,
    min_chars: int = 200,
    max_chars: int = 1800,
) -> list[str]:
    """
    Split text into semantically coherent chunks using a self-calibrating
    embedding-gradient detector.

    Algorithm:
    1. Sentence-tokenise the text.
    2. For each candidate boundary i, compute cosine sim between the `window`
       sentences BEFORE and AFTER it.
    3. Treat the bottom `percentile` of those sim scores as topic-shift breakpoints.
       (self-calibrating: no fixed threshold needed)
    4. Merge micro-chunks < min_chars; hard cap at max_chars.

    Returns a list of chunk strings (no metadata).
    """
    import re

    # ── 1. Sentence tokenise ──────────────────────────────────────────────────
    sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", text) if s.strip()]
    if len(sentences) <= 2 * window + 1:
        # Too short to find meaningful boundaries
        return [text] if text.strip() else []

    # ── 2. Batch compute before/after cosine sims at each candidate boundary ──
    windows_before: list[str] = []
    windows_after: list[str] = []
    indices: list[int] = []
    for i in range(window, len(sentences) - window):
        before = " ".join(sentences[max(0, i - window): i])
        after  = " ".join(sentences[i: i + window])
        windows_before.append(before)
        windows_after.append(after)
        indices.append(i)

    if not indices:
        return [text]

    all_texts = windows_before + windows_after
    try:
        if hasattr(embedder, "embed_documents"):
            all_vecs = embedder.embed_documents(all_texts)
        else:
            all_vecs = [embedder.embed_query(t) for t in all_texts]
    except Exception:
        all_vecs = [embedder.embed_query(t) for t in all_texts]

    n_pairs = len(indices)
    before_vecs = all_vecs[:n_pairs]
    after_vecs = all_vecs[n_pairs:]

    sims: list[tuple[int, float]] = []
    for idx, (b_vec, a_vec) in enumerate(zip(before_vecs, after_vecs)):
        sim = _cosine_sim(b_vec, a_vec)
        sims.append((indices[idx], sim))

    if not sims:
        return [text]


    # ── 3. Percentile-based breakpoint detection ──────────────────────────────
    sim_values = sorted(s for _, s in sims)
    n = len(sim_values)
    cutoff_idx = max(0, int(n * percentile / 100.0) - 1)
    cutoff = sim_values[cutoff_idx]

    breakpoints = [i for i, s in sims if s <= cutoff]

    # Deduplicate: among adjacent breakpoints, keep the one with lowest sim (sharpest boundary)
    sim_lookup = dict(sims)
    clean_bps: list[int] = []
    cluster: list[int] = []
    for bp in sorted(set(breakpoints)):
        if cluster and bp - cluster[-1] <= window:
            cluster.append(bp)
        else:
            if cluster:
                # Keep the sharpest boundary in this cluster
                best = min(cluster, key=lambda x: sim_lookup.get(x, 1.0))
                clean_bps.append(best)
            cluster = [bp]
    if cluster:
        best = min(cluster, key=lambda x: sim_lookup.get(x, 1.0))
        clean_bps.append(best)
    clean_bps.sort()

    if not clean_bps:
        return [text]

    # ── 4. Build chunk strings ────────────────────────────────────────────────
    chunks: list[str] = []
    prev = 0
    for bp in clean_bps:
        chunk_text = " ".join(sentences[prev:bp]).strip()
        if chunk_text:
            chunks.append(chunk_text)
        prev = bp
    tail = " ".join(sentences[prev:]).strip()
    if tail:
        chunks.append(tail)

    if not chunks:
        return [text]

    # ── 5. Merge only micro-chunks / hard-cap giant chunks ────────────────────
    merged: list[str] = []
    buf = ""
    for chunk in chunks:
        if not buf:
            buf = chunk
            continue
        # Only merge if the current buffer is a micro-chunk (below min threshold)
        if len(buf) < min_chars:
            candidate = (buf + " " + chunk).strip()
            if len(candidate) <= max_chars:
                buf = candidate
            else:
                merged.append(buf)
                buf = chunk
        else:
            # buf is a proper-sized chunk — emit it
            merged.append(buf)
            buf = chunk
    if buf:
        # Hard-cap: if final buf exceeds max, character-split it
        if len(buf) > max_chars:
            from langchain_text_splitters import RecursiveCharacterTextSplitter
            splitter = RecursiveCharacterTextSplitter(chunk_size=max_chars, chunk_overlap=80)
            merged.extend([c.page_content if hasattr(c, 'page_content') else c
                           for c in splitter.create_documents([buf])])
        else:
            merged.append(buf)

    return merged if merged else [text]


def semantic_split_documents(
    sections: list[Document],
    embedder,
    window: int = 3,
    percentile: float = 25.0,
    min_chars: int = 200,
    max_chars: int = 1800,
) -> list[Document]:
    """
    Apply semantic_chunk to each Document section, returning Document objects
    with metadata preserved (breadcrumb, headings, source).
    Falls back gracefully if fewer than (2*window+1) sentences detected.
    """
    result: list[Document] = []
    for section in sections:
        text = section.page_content.strip()
        if not text:
            continue

        # Sections shorter than min_chars go straight through
        if len(text) < min_chars:
            result.append(ensure_chunk_has_headings(section))
            continue

        try:
            chunk_texts = semantic_chunk(
                text, embedder,
                window=window, percentile=percentile,
                min_chars=min_chars, max_chars=max_chars,
            )
        except Exception:
            # Graceful degradation
            chunk_texts = [text]

        for ct in chunk_texts:
            doc = Document(page_content=ct, metadata=dict(section.metadata))
            result.append(ensure_chunk_has_headings(doc))

    return result



def is_url(source: str) -> bool:
    return urlparse(source).scheme in ("http", "https")


# ---------------------------------------------------------------------------
# RapidOCR Engine (Lazy Initialized)
# ---------------------------------------------------------------------------

_ocr_engine = None

def get_ocr_engine():
    global _ocr_engine
    if _ocr_engine is None:
        from rapidocr_onnxruntime import RapidOCR
        _ocr_engine = RapidOCR()
    return _ocr_engine


def extract_text_from_image(image_input) -> str:
    """Extracts text lines from an image path, PIL Image, or numpy array via RapidOCR."""
    import numpy as np
    from PIL import Image

    if isinstance(image_input, (str, Path)):
        img = Image.open(image_input).convert("RGB")
        img_np = np.array(img)
    elif isinstance(image_input, Image.Image):
        img_np = np.array(image_input.convert("RGB"))
    else:
        img_np = image_input

    engine = get_ocr_engine()
    result, _ = engine(img_np)
    if not result:
        return ""

    lines = [item[1] for item in result if item and len(item) > 1]
    return "\n".join(lines).strip()


# ---------------------------------------------------------------------------
# 1. PDF Loader (with Automatic Scanned PDF OCR Fallback)
# ---------------------------------------------------------------------------

def load_pdf(path: str, chunk_size: int, chunk_overlap: int) -> list[Document]:
    from langchain_community.document_loaders import PyPDFLoader
    import pypdf
    import io
    from PIL import Image

    docs = PyPDFLoader(path).load()
    total_text = "".join([d.page_content for d in docs]).strip()

    # If PDF has native selectable text, use fast digital extraction
    if len(total_text) >= 50:
        for d in docs:
            d.metadata["source"] = path
        return _sub_chunk(docs, chunk_size, chunk_overlap)

    # Scanned PDF detected -> Fallback to RapidOCR on page images
    print(f"  [OCR Fallback] Scanned PDF detected for '{Path(path).name}'. Running OCR on pages...")
    reader = pypdf.PdfReader(path)
    ocr_docs = []

    for i, page in enumerate(reader.pages, start=1):
        page_text_lines = []
        for img_obj in page.images:
            try:
                pil_img = Image.open(io.BytesIO(img_obj.data))
                extracted = extract_text_from_image(pil_img)
                if extracted:
                    page_text_lines.append(extracted)
            except Exception as e:
                print(f"    OCR failed on page {i} image: {e}")

        if page_text_lines:
            page_content = f"## Page {i} (Scanned Document)\n" + "\n".join(page_text_lines)
            ocr_docs.append(Document(page_content=page_content, metadata={"source": path, "h1": Path(path).stem, "h2": f"Page {i}"}))

    if not ocr_docs:
        # Fallback to whatever tiny text was found or empty placeholder
        ocr_docs = [Document(page_content=total_text or "[Empty Scanned PDF]", metadata={"source": path})]

    return _sub_chunk(ocr_docs, chunk_size, chunk_overlap)


# ---------------------------------------------------------------------------
# 2. Image Loader (.png, .jpg, .jpeg, .webp, .bmp, .tiff)
# ---------------------------------------------------------------------------

def load_image(path: str, chunk_size: int, chunk_overlap: int) -> list[Document]:
    filename = Path(path).name
    extracted = extract_text_from_image(path)
    content = f"# Image Document: {filename}\n\n{extracted or '[No text detected in image]'}"
    doc = Document(page_content=content, metadata={"source": path, "h1": filename, "type": "image_ocr"})
    return _sub_chunk([doc], chunk_size, chunk_overlap)


# ---------------------------------------------------------------------------
# 2. Microsoft Word Loader (.docx)
# ---------------------------------------------------------------------------

def load_docx(path: str, chunk_size: int, chunk_overlap: int) -> list[Document]:
    import docx

    doc = docx.Document(path)
    sections = []
    current_h1 = Path(path).stem
    current_h2 = ""
    current_lines = []

    def flush():
        if current_lines:
            meta = {"source": path, "h1": current_h1}
            if current_h2:
                meta["h2"] = current_h2
            sections.append(Document(page_content="\n".join(current_lines).strip(), metadata=meta))

    for p in doc.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        style = p.style.name.lower() if p.style else ""
        if "heading 1" in style:
            flush()
            current_h1 = text
            current_h2 = ""
            current_lines = [f"# {text}"]
        elif "heading 2" in style:
            flush()
            current_h2 = text
            current_lines = [f"## {text}"]
        elif "heading 3" in style:
            current_lines.append(f"### {text}")
        else:
            current_lines.append(text)

    # Extract tables in document as markdown tables
    for table in doc.tables:
        rows_text = []
        for i, row in enumerate(table.rows):
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            rows_text.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows_text.append("| " + " | ".join(["---"] * len(cells)) + " |")
        if rows_text:
            current_lines.append("\n" + "\n".join(rows_text) + "\n")

    flush()
    if not sections:
        sections = [Document(page_content="[Empty Word Document]", metadata={"source": path})]
    return _sub_chunk(sections, chunk_size, chunk_overlap)


# ---------------------------------------------------------------------------
# 3. Microsoft PowerPoint Loader (.pptx)
# ---------------------------------------------------------------------------

def load_pptx(path: str, chunk_size: int, chunk_overlap: int) -> list[Document]:
    import pptx

    prs = pptx.Presentation(path)
    sections = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_title = f"Slide {i}"
        slide_text = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        if not slide_text and hasattr(shape, "is_placeholder") and shape.is_placeholder:
                            slide_title = f"Slide {i}: {t}"
                        slide_text.append(t)

        content = f"## {slide_title}\n" + "\n".join(slide_text)
        sections.append(Document(page_content=content, metadata={"source": path, "h1": Path(path).stem, "h2": slide_title}))

    if not sections:
        sections = [Document(page_content="[Empty Presentation]", metadata={"source": path})]
    return _sub_chunk(sections, chunk_size, chunk_overlap)


# ---------------------------------------------------------------------------
# 4. Tabular Data: Excel (.xlsx) & CSV/TSV
# ---------------------------------------------------------------------------

def load_excel(path: str, chunk_size: int, chunk_overlap: int) -> list[Document]:
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True)
    sections = []

    for sheetname in wb.sheetnames:
        sheet = wb[sheetname]
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue

        headers = [str(c or "").strip() for c in rows[0]]
        table_lines = [f"## Sheet: {sheetname}", "| " + " | ".join(headers) + " |", "| " + " | ".join(["---"] * len(headers)) + " |"]

        for row in rows[1:1000]:  # Limit to 1000 rows per sheet
            if any(row):
                cells = [str(c or "").strip().replace("\n", " ") for c in row]
                table_lines.append("| " + " | ".join(cells) + " |")

        sections.append(Document(page_content="\n".join(table_lines), metadata={"source": path, "h1": Path(path).stem, "h2": f"Sheet: {sheetname}"}))

    if not sections:
        sections = [Document(page_content="[Empty Spreadsheet]", metadata={"source": path})]
    return _sub_chunk(sections, chunk_size, chunk_overlap)


def load_csv(path: str, chunk_size: int, chunk_overlap: int, delimiter: str = ",") -> list[Document]:
    sections = []
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f, delimiter=delimiter)
        rows = list(reader)

    if rows:
        headers = rows[0]
        table_lines = [f"## File: {Path(path).name}", "| " + " | ".join(headers) + " |", "| " + " | ".join(["---"] * len(headers)) + " |"]
        for r in rows[1:1500]:
            if r:
                table_lines.append("| " + " | ".join([c.strip().replace("\n", " ") for c in r]) + " |")
        sections.append(Document(page_content="\n".join(table_lines), metadata={"source": path, "h1": Path(path).name}))
    else:
        sections.append(Document(page_content="[Empty CSV File]", metadata={"source": path}))

    return _sub_chunk(sections, chunk_size, chunk_overlap)


# ---------------------------------------------------------------------------
# 5. Markdown & Code / Config Loaders
# ---------------------------------------------------------------------------

def load_markdown(path: str, chunk_size: int, chunk_overlap: int) -> list[Document]:
    text = Path(path).read_text(encoding="utf-8", errors="replace")
    headers_to_split_on = [("#", "h1"), ("##", "h2"), ("###", "h3")]
    md_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
    sections = md_splitter.split_text(text)

    for s in sections:
        s.metadata["source"] = path

    formatted = [format_section(s, source=path) for s in sections]
    return _sub_chunk(formatted, chunk_size, chunk_overlap)


def load_code_or_text(path: str, chunk_size: int, chunk_overlap: int) -> list[Document]:
    text = Path(path).read_text(encoding="utf-8", errors="replace")
    ext = Path(path).suffix.lower()
    filename = Path(path).name

    # Wrap code files with markdown code blocks for syntactic context
    if ext in (".py", ".js", ".ts", ".tsx", ".jsx", ".cpp", ".c", ".h", ".java", ".go", ".rs", ".sql", ".sh", ".json", ".yaml", ".yml", ".toml"):
        lang = ext.lstrip(".")
        if lang in ("yml", "yaml"): lang = "yaml"
        if lang in ("ts", "tsx"): lang = "typescript"
        if lang in ("js", "jsx"): lang = "javascript"
        if lang in ("py",): lang = "python"
        
        formatted_content = f"## Code File: {filename}\n```{lang}\n{text}\n```"
    else:
        formatted_content = f"## Document: {filename}\n{text}"

    doc = Document(page_content=formatted_content, metadata={"source": path, "h1": filename})
    return _sub_chunk([doc], chunk_size, chunk_overlap)


# ---------------------------------------------------------------------------
# 6. YouTube Transcripts Loader
# ---------------------------------------------------------------------------

def is_youtube_url(url: str) -> bool:
    parsed = urlparse(url)
    return "youtube.com" in parsed.netloc or "youtu.be" in parsed.netloc


def extract_youtube_video_id(url: str) -> Optional[str]:
    parsed = urlparse(url)
    if "youtu.be" in parsed.netloc:
        return parsed.path.lstrip("/")
    if "youtube.com" in parsed.netloc:
        qs = parse_qs(parsed.query)
        if "v" in qs:
            return qs["v"][0]
        if parsed.path.startswith("/shorts/"):
            return parsed.path.split("/shorts/")[1].split("/")[0]
        if parsed.path.startswith("/embed/"):
            return parsed.path.split("/embed/")[1].split("/")[0]
    return None


def load_youtube_transcript(url: str, chunk_size: int, chunk_overlap: int) -> list[Document]:
    video_id = extract_youtube_video_id(url)
    if not video_id:
        raise ValueError(f"Could not extract YouTube video ID from URL: {url}")

    # 1. Fetch Video Title via YouTube oEmbed API
    title = f"YouTube Video ({video_id})"
    try:
        oembed_res = requests.get(f"https://www.youtube.com/oembed?url=https://www.youtube.com/watch?v={video_id}&format=json", timeout=5)
        if oembed_res.status_code == 200:
            title = oembed_res.json().get("title", title)
    except Exception:
        pass

    # 2. Fetch Transcript
    from youtube_transcript_api import YouTubeTranscriptApi

    try:
        transcript_data = YouTubeTranscriptApi.get_transcript(video_id)
    except Exception as e:
        # Try fetching list of available languages
        try:
            transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
            transcript = transcript_list.find_transcript(['en', 'en-US', 'auto'])
            transcript_data = transcript.fetch()
        except Exception as e2:
            raise ValueError(f"Could not retrieve transcripts for YouTube video '{title}': {e2}")

    # Format transcript with timestamps
    sections = []
    current_lines = [f"# YouTube Video: {title}", f"Source URL: {url}\n"]
    current_time_str = "00:00"

    for entry in transcript_data:
        start_seconds = int(entry.get("start", 0))
        mins, secs = divmod(start_seconds, 60)
        hours, mins = divmod(mins, 60)
        time_tag = f"[{hours:02d}:{mins:02d}:{secs:02d}]" if hours > 0 else f"[{mins:02d}:{secs:02d}]"
        
        text = entry.get("text", "").replace("\n", " ").strip()
        if text:
            current_lines.append(f"{time_tag} {text}")

    full_transcript_text = "\n".join(current_lines)
    doc = Document(page_content=full_transcript_text, metadata={"source": url, "h1": title, "type": "youtube"})
    return _sub_chunk([doc], chunk_size, chunk_overlap)


# ---------------------------------------------------------------------------
# 7. GitHub Raw & Web URLs Loader
# ---------------------------------------------------------------------------

def load_github_url(url: str, chunk_size: int, chunk_overlap: int, user_agent: str) -> list[Document]:
    # Convert github.com/owner/repo/blob/branch/file to raw.githubusercontent.com
    raw_url = url.replace("github.com", "raw.githubusercontent.com").replace("/blob/", "/")
    res = requests.get(raw_url, timeout=15, headers={"User-Agent": user_agent})
    res.raise_for_status()

    filename = Path(urlparse(url).path).name or "GitHub Document"
    text = res.text
    if url.endswith(".md") or url.endswith(".markdown"):
        headers_to_split_on = [("#", "h1"), ("##", "h2"), ("###", "h3")]
        md_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
        sections = md_splitter.split_text(text)
        for s in sections:
            s.metadata["source"] = url
        return _sub_chunk([format_section(s, source=url) for s in sections], chunk_size, chunk_overlap)

    doc = Document(page_content=f"## File: {filename}\n```\n{text}\n```", metadata={"source": url, "h1": filename})
    return _sub_chunk([doc], chunk_size, chunk_overlap)


def load_header_sections(url: str, user_agent: str = DEFAULT_USER_AGENT) -> list[Document]:
    response = requests.get(url, timeout=30, headers={"User-Agent": user_agent})
    response.raise_for_status()

    soup = BeautifulSoup(response.text, "html.parser")
    for noisy in soup.select("nav, .toc, .table-of-contents, #table-of-contents, script, style, footer"):
        noisy.decompose()

    page_title = (soup.title.string if soup.title else "").strip() or urlparse(url).netloc
    article = soup.select_one("article") or soup.select_one(".post-content") or soup.select_one("main") or soup.body

    sections = []
    current_metadata = {"h1": page_title}
    current_lines = [f"# {page_title}"]

    def flush_section():
        if current_lines:
            sections.append(
                format_section(
                    Document(
                        page_content="\n".join(current_lines).strip(),
                        metadata=current_metadata.copy(),
                    ),
                    source=url,
                )
            )

    for element in article.find_all(["h1", "h2", "h3", "p", "li", "pre", "table"]):
        if element.name in {"h1", "h2", "h3"}:
            flush_section()
            current_lines = []

            level = int(element.name[1])
            heading = clean_heading(element.get_text(" ", strip=True))
            if level == 1:
                current_metadata = {"h1": heading}
            elif level == 2:
                current_metadata = {"h1": current_metadata.get("h1", page_title), "h2": heading}
            else:
                current_metadata = {
                    "h1": current_metadata.get("h1", page_title),
                    "h2": current_metadata.get("h2", ""),
                    "h3": heading,
                }
            continue

        text = element.get_text(" ", strip=True)
        if text:
            current_lines.append(text)

    flush_section()
    return sections


def load_url(url: str, chunk_size: int, chunk_overlap: int, user_agent: str) -> list[Document]:
    if is_youtube_url(url):
        return load_youtube_transcript(url, chunk_size, chunk_overlap)

    if "github.com" in url and ("/blob/" in url or "/raw/" in url):
        return load_github_url(url, chunk_size, chunk_overlap, user_agent)

    sections = load_header_sections(url, user_agent=user_agent)
    return _sub_chunk(sections, chunk_size, chunk_overlap)


# ---------------------------------------------------------------------------
# Universal Dispatcher
# ---------------------------------------------------------------------------

def load_and_split_source(
    source: str,
    chunk_size: int | None = None,
    chunk_overlap: int | None = None,
    user_agent: str | None = None,
) -> list[Document]:
    size = chunk_size or int(os.getenv("CHUNK_SIZE", DEFAULT_CHUNK_SIZE))
    overlap = chunk_overlap or int(os.getenv("CHUNK_OVERLAP", DEFAULT_CHUNK_OVERLAP))
    agent = user_agent or os.getenv("USER_AGENT", DEFAULT_USER_AGENT)

    # 1. URL Dispatching
    if is_url(source):
        return load_url(source, size, overlap, agent)

    # 2. File Extension Dispatching
    suffix = Path(source).suffix.lower()

    if suffix == ".pdf":
        return load_pdf(source, size, overlap)
    
    if suffix in (".docx", ".doc"):
        return load_docx(source, size, overlap)

    if suffix in (".pptx", ".ppt"):
        return load_pptx(source, size, overlap)

    if suffix in (".xlsx", ".xls"):
        return load_excel(source, size, overlap)

    if suffix == ".csv":
        return load_csv(source, size, overlap, delimiter=",")

    if suffix == ".tsv":
        return load_csv(source, size, overlap, delimiter="\t")

    if suffix in (".md", ".markdown"):
        return load_markdown(source, size, overlap)

    # Images (PNG, JPG, JPEG, WEBP, BMP, TIFF) via RapidOCR
    if suffix in (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff"):
        return load_image(source, size, overlap)

    # Code, Config, Subtitles, Plain Text
    if suffix in (
        ".txt", ".py", ".js", ".ts", ".tsx", ".jsx", ".json", ".yaml", ".yml",
        ".toml", ".ini", ".sql", ".html", ".htm", ".css", ".cpp", ".c", ".h",
        ".java", ".go", ".rs", ".sh", ".bash", ".zsh", ".srt", ".vtt", ""
    ):
        return load_code_or_text(source, size, overlap)

    raise ValueError(
        f"Unsupported file format '{suffix}'. Supported: .pdf, .png, .jpg, .docx, .pptx, .xlsx, .csv, .md, .txt, .py, .js, .ts, .json, .yaml, YouTube URLs, and Web URLs."
    )


def ingest_document_structure_aware(
    source: str,
    original_filename: Optional[str] = None,
    user_id: Optional[str] = None,
    target_chunk_size: int = 1200,
    chunk_overlap: int = 150,
    child_chunk_size: int = 350,
    child_overlap: int = 50,
):
    """
    Parses any document source via UnifiedDocumentParser into a DocumentAST,
    then executes element-aware semantic chunking into (parent_docs, child_docs)
    and returns comprehensive lineage telemetry.
    """
    import time
    from app.document_intelligence.parser import get_document_parser
    from app.document_intelligence.chunker import StructureAwareChunker
    from app.retrieval.contextual import get_contextual_engine

    t0 = time.time()
    parser = get_document_parser()
    doc_ast, parser_name, parser_version = parser.parse(source, original_filename=original_filename)

    chunker = StructureAwareChunker(
        target_chunk_size=target_chunk_size,
        chunk_overlap=chunk_overlap,
        child_chunk_size=child_chunk_size,
        child_overlap=child_overlap,
    )
    parent_chunks, child_chunks = chunker.chunk_document(doc_ast)

    # Deduplication & Boilerplate Stripping (SHA-256 + SimHash)
    from app.document_intelligence.dedup import get_deduplicator
    deduplicator = get_deduplicator()
    child_chunks, dedup_removed_count = deduplicator.deduplicate_chunks(child_chunks)


    # Contextual Retrieval Enrichment
    contextual_engine = get_contextual_engine()
    doc_full_text = doc_ast.to_markdown()
    child_chunks = contextual_engine.enrich_chunks(
        chunks=child_chunks,
        doc_text=doc_full_text,
        doc_title=original_filename or doc_ast.filename,
    )

    # Hierarchical Document Summary
    from app.document_intelligence.summarizer import get_hierarchical_summarizer
    summarizer = get_hierarchical_summarizer()
    summary_chunk = summarizer.generate_document_summary(doc_ast)
    if summary_chunk:
        child_chunks.append(summary_chunk)

    # Convert StructuredChunks to LangChain Document objects with lineage & metadata
    parent_docs = []

    for p in parent_chunks:
        meta = {
            **p.metadata,
            "chunk_id": p.id,
            "is_parent": True,
            "content_type": p.content_type,
            "raw_content": p.raw_content,
            "page": p.page_number,
            "h1": p.heading,
            "h2": p.section,
            "section_path": p.section_path,
            "user_id": user_id,
        }
        parent_docs.append(Document(page_content=p.content, metadata=meta))

    child_docs = []
    for c in child_chunks:
        meta = {
            **c.metadata,
            "chunk_id": c.id,
            "parent_id": c.parent_id,
            "is_parent": False,
            "content_type": c.content_type,
            "raw_content": c.raw_content,
            "contextual_content": c.contextual_content,
            "page": c.page_number,
            "h1": c.heading,
            "h2": c.section,
            "section_path": c.section_path,
            "user_id": user_id,
        }
        child_docs.append(Document(page_content=c.content, metadata=meta))

    elapsed_ms = int((time.time() - t0) * 1000)
    lineage_info = {
        "parser_name": parser_name,
        "parser_version": parser_version,
        "chunker_version": chunker.version,
        "contextualization_model": "contextual_engine_v1",
        "processing_time_ms": elapsed_ms,
        "chunk_count": len(child_docs),
        "parent_count": len(parent_docs),
        "table_count": len(doc_ast.all_tables()),
        "figure_count": len(doc_ast.all_figures()),
        "ocr_page_count": doc_ast.ocr_page_count(),
        "dedup_removed_count": dedup_removed_count,
    }


    return doc_ast, parent_docs, child_docs, lineage_info



def load_and_split_sources(sources: list[str], **kwargs) -> list[Document]:
    all_splits = []
    for source in sources:
        print(f"  Loading: {source}")
        splits = load_and_split_source(source, **kwargs)
        print(f"    -> {len(splits)} chunks")
        all_splits.extend(splits)
    return all_splits


def split_article(
    url: str | None = None,
    chunk_size: int | None = None,
    chunk_overlap: int | None = None,
    user_agent: str | None = None,
) -> list[Document]:
    article_url = url or os.getenv("ARTICLE_URL", DEFAULT_ARTICLE_URL)
    return load_and_split_source(
        article_url,
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        user_agent=user_agent,
    )